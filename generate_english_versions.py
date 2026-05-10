from __future__ import annotations

import json
from pathlib import Path

import nbformat as nbf
import pandas as pd
import plotly.express as px
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptInches
from pptx.util import Pt


BASE_DIR = Path(__file__).resolve().parent
DATA_FILE = BASE_DIR / "Marketing_Campaign_Data.csv"
DELIVERABLES_DIR = BASE_DIR / "deliverables"
CHARTS_DIR = DELIVERABLES_DIR / "charts"


def load_data() -> pd.DataFrame:
    df = pd.read_csv(DATA_FILE)
    return df.rename(
        columns={
            "Interaction ID": "interaction_id",
            "Campaign Type": "campaign",
            "Channel": "channel",
            "Customer Type": "customer_type",
            "Converted (1=yes, 0=no)": "converted",
            "Time on Site (seconds)": "time_on_site",
            "Sales ($)": "sales",
        }
    )


def compute_tables(df: pd.DataFrame) -> tuple[dict, dict]:
    kpis = {
        "interactions": int(len(df)),
        "total_revenue": float(df["sales"].sum()),
        "conversion_rate": float(df["converted"].mean()),
        "avg_time": float(df["time_on_site"].mean()),
        "avg_ticket": float(df.loc[df["converted"] == 1, "sales"].mean()),
        "new_share": float((df["customer_type"] == "New").mean()),
    }

    tables = {
        "by_channel": (
            df.groupby("channel", as_index=False)
            .agg(interactions=("interaction_id", "count"), conversion_rate=("converted", "mean"), revenue=("sales", "sum"))
            .sort_values("revenue", ascending=False)
        ),
        "by_campaign": (
            df.groupby("campaign", as_index=False)
            .agg(interactions=("interaction_id", "count"), conversion_rate=("converted", "mean"), revenue=("sales", "sum"))
            .sort_values("revenue", ascending=False)
        ),
        "by_combo": (
            df.groupby(["campaign", "channel"], as_index=False)
            .agg(interactions=("interaction_id", "count"), conversion_rate=("converted", "mean"), revenue=("sales", "sum"))
            .sort_values("revenue", ascending=False)
        ),
    }
    tables["by_combo"]["combo"] = tables["by_combo"]["campaign"] + " - " + tables["by_combo"]["channel"]

    return kpis, tables


def build_summary_en(kpis: dict, tables: dict) -> Path:
    top_combo = tables["by_combo"].iloc[0]
    output = DELIVERABLES_DIR / "marketing_campaign_analysis_summary_en.md"
    text = f"""
# Marketing Campaign Analysis Summary

## Dataset
- Rows: {kpis['interactions']:,}
- Campaigns: A and B
- Channels: Email, Instagram, Web Banner

## Executive KPIs
- Total revenue: ${kpis['total_revenue']:,.2f}
- Conversion rate: {kpis['conversion_rate'] * 100:.2f}%
- Average time on site: {kpis['avg_time']:.2f} seconds
- Average ticket (converted only): ${kpis['avg_ticket']:.2f}

## Top Insight
- Highest revenue combination: {top_combo['campaign']} + {top_combo['channel']} (${top_combo['revenue']:,.2f})

## Recommendations
1. Prioritize budget toward B + Email for efficient scale.
2. Run A/B tests by channel to improve conversion.
3. Segment strategies for New vs Existing customers.
4. Review KPI performance weekly and rebalance spend.
""".strip()
    output.write_text(text + "\n", encoding="utf-8")
    return output


def build_executive_dashboard_en(df: pd.DataFrame, kpis: dict, tables: dict) -> Path:
    output = DELIVERABLES_DIR / "marketing_campaign_executive_dashboard_en.html"

    by_channel = tables["by_channel"].copy()
    by_campaign = tables["by_campaign"].copy()
    by_combo = tables["by_combo"].copy()
    by_campaign["conversion_rate"] *= 100

    fig1 = px.bar(by_channel, x="channel", y="revenue", color="channel", title="Revenue by Channel", text_auto=".2s")
    fig1.update_layout(showlegend=False)

    fig2 = px.bar(by_campaign, x="campaign", y="conversion_rate", color="campaign", title="Conversion Rate (%) by Campaign", text_auto=".2f")
    fig2.update_layout(showlegend=False, yaxis_title="Conversion rate (%)")

    fig3 = px.bar(by_combo, x="combo", y="revenue", color="campaign", title="Revenue by Campaign + Channel Combination")
    fig3.update_layout(xaxis_title="Combination", yaxis_title="Revenue ($)")

    converted = df[df["converted"] == 1]
    fig4 = px.box(converted, x="channel", y="sales", color="campaign", title="Ticket Distribution (Converted Interactions)")

    html = f"""
<!DOCTYPE html>
<html lang=\"en\">
<head>
  <meta charset=\"UTF-8\" />
  <meta name=\"viewport\" content=\"width=device-width, initial-scale=1.0\" />
  <title>Marketing Campaign Executive Dashboard (EN)</title>
  <script src=\"https://cdn.plot.ly/plotly-2.35.2.min.js\"></script>
  <style>
    body {{ margin:0; font-family: Georgia, 'Times New Roman', serif; background:#f4f1ea; color:#1f2a37; }}
    .wrap {{ max-width:1280px; margin:0 auto; padding:24px; }}
    .cards {{ display:grid; grid-template-columns:repeat(auto-fit,minmax(180px,1fr)); gap:12px; margin-bottom:18px; }}
    .card {{ background:#fff; border-radius:12px; padding:12px; box-shadow:0 6px 18px rgba(0,0,0,.08); }}
    .plot {{ background:#fff; border-radius:12px; padding:8px; box-shadow:0 6px 18px rgba(0,0,0,.08); margin-bottom:12px; }}
    h1 {{ margin:0 0 6px; }}
  </style>
</head>
<body>
  <div class=\"wrap\">
    <h1>Marketing Campaign Executive Dashboard</h1>
    <p>Dataset: 100,000 interactions across Campaign A and B with Email, Instagram, and Web Banner channels.</p>
    <div class=\"cards\">
      <div class=\"card\"><b>Total Interactions</b><div>{kpis['interactions']:,}</div></div>
      <div class=\"card\"><b>Total Revenue</b><div>${kpis['total_revenue']:,.2f}</div></div>
      <div class=\"card\"><b>Conversion Rate</b><div>{kpis['conversion_rate']*100:.2f}%</div></div>
      <div class=\"card\"><b>Avg Time on Site</b><div>{kpis['avg_time']:.1f}s</div></div>
      <div class=\"card\"><b>Avg Ticket</b><div>${kpis['avg_ticket']:.2f}</div></div>
      <div class=\"card\"><b>New Customer Share</b><div>{kpis['new_share']*100:.2f}%</div></div>
    </div>
    <div class=\"plot\" id=\"plot1\"></div>
    <div class=\"plot\" id=\"plot2\"></div>
    <div class=\"plot\" id=\"plot3\"></div>
    <div class=\"plot\" id=\"plot4\"></div>
  </div>
  <script>
    Plotly.newPlot('plot1', {fig1.to_json()}['data'], {fig1.to_json()}['layout'], {{responsive: true}});
    Plotly.newPlot('plot2', {fig2.to_json()}['data'], {fig2.to_json()}['layout'], {{responsive: true}});
    Plotly.newPlot('plot3', {fig3.to_json()}['data'], {fig3.to_json()}['layout'], {{responsive: true}});
    Plotly.newPlot('plot4', {fig4.to_json()}['data'], {fig4.to_json()}['layout'], {{responsive: true}});
  </script>
</body>
</html>
"""
    output.write_text(html, encoding="utf-8")
    return output


def build_interactive_dashboard_en(df: pd.DataFrame) -> Path:
    output = DELIVERABLES_DIR / "marketing_campaign_dashboard_interactive_filters_en.html"
    data = json.dumps(df[["campaign", "channel", "customer_type", "converted", "time_on_site", "sales"]].to_dict(orient="records"))

    html = f"""
<!DOCTYPE html>
<html lang=\"en\">
<head>
  <meta charset=\"UTF-8\" />
  <meta name=\"viewport\" content=\"width=device-width, initial-scale=1.0\" />
  <title>Interactive Dashboard with Filters (EN)</title>
  <script src=\"https://cdn.plot.ly/plotly-2.35.2.min.js\"></script>
</head>
<body style=\"font-family:Cambria,Georgia,serif;background:#f6efe0;margin:0;\">
<div style=\"max-width:1200px;margin:0 auto;padding:20px;\">
  <h1>Interactive Dashboard with Filters</h1>
  <p>Filter by campaign, channel, and customer type.</p>
  <div style=\"display:flex;gap:10px;flex-wrap:wrap;\">
    <select id=\"campaign\"></select>
    <select id=\"channel\"></select>
    <select id=\"customer\"></select>
  </div>
  <div id=\"cards\" style=\"display:grid;grid-template-columns:repeat(auto-fit,minmax(170px,1fr));gap:8px;margin:12px 0;\"></div>
  <div id=\"p1\" style=\"height:360px;\"></div>
  <div id=\"p2\" style=\"height:360px;\"></div>
  <div id=\"p3\" style=\"height:360px;\"></div>
</div>
<script>
const raw = {data};
const c = document.getElementById('campaign');
const ch = document.getElementById('channel');
const ct = document.getElementById('customer');

function uniq(k) {{ return [...new Set(raw.map(r => r[k]))].sort(); }}
function fill(sel, arr) {{
  sel.innerHTML='';
  const a=document.createElement('option'); a.value='All'; a.textContent='All'; sel.appendChild(a);
  arr.forEach(v=>{{ const o=document.createElement('option'); o.value=v; o.textContent=v; sel.appendChild(o); }});
}}
function filt() {{
  return raw.filter(r => (c.value==='All'||r.campaign===c.value) && (ch.value==='All'||r.channel===ch.value) && (ct.value==='All'||r.customer_type===ct.value));
}}
function sumBy(data,key,val) {{
  const m=new Map(); data.forEach(r=>m.set(r[key],(m.get(r[key])||0)+r[val]));
  return [...m.entries()].map(([k,v])=>({{k,v}}));
}}
function meanBy(data,key,val) {{
  const m=new Map(); data.forEach(r=>{{const e=m.get(r[key])||{{s:0,n:0}};e.s+=r[val];e.n+=1;m.set(r[key],e);}});
  return [...m.entries()].map(([k,o])=>({{k,v:o.n?o.s/o.n:0}}));
}}
function render() {{
  const d=filt();
  const n=d.length; const rev=d.reduce((a,b)=>a+b.sales,0); const cr=n?d.reduce((a,b)=>a+b.converted,0)/n:0;
  const avg=n?d.reduce((a,b)=>a+b.time_on_site,0)/n:0; const conv=d.filter(x=>x.converted===1);
  const ticket=conv.length?conv.reduce((a,b)=>a+b.sales,0)/conv.length:0;
  document.getElementById('cards').innerHTML = `
    <div><b>Interactions</b><div>${{n.toLocaleString()}}</div></div>
    <div><b>Revenue</b><div>$${{rev.toLocaleString(undefined,{{maximumFractionDigits:2}})}}</div></div>
    <div><b>Conversion</b><div>${{(cr*100).toFixed(2)}}%</div></div>
    <div><b>Avg Time</b><div>${{avg.toFixed(1)}}s</div></div>
    <div><b>Avg Ticket</b><div>$${{ticket.toFixed(2)}}</div></div>`;
  const rch=sumBy(d,'channel','sales').sort((a,b)=>b.v-a.v);
  const rc=meanBy(d,'campaign','converted').sort((a,b)=>b.v-a.v);
  const comboMap=new Map(); d.forEach(r=>{{const key=`${{r.campaign}} - ${{r.channel}}`; comboMap.set(key,(comboMap.get(key)||0)+r.sales);}});
  const rcombo=[...comboMap.entries()].map(([k,v])=>({{k,v}})).sort((a,b)=>b.v-a.v);
  Plotly.react('p1',[{{type:'bar',x:rch.map(x=>x.k),y:rch.map(x=>x.v)}}],{{title:'Revenue by Channel'}},{{responsive:true}});
  Plotly.react('p2',[{{type:'bar',x:rc.map(x=>x.k),y:rc.map(x=>x.v*100)}}],{{title:'Conversion by Campaign (%)'}},{{responsive:true}});
  Plotly.react('p3',[{{type:'bar',x:rcombo.map(x=>x.k),y:rcombo.map(x=>x.v)}}],{{title:'Revenue by Campaign-Channel Combination'}},{{responsive:true}});
}}
fill(c,uniq('campaign')); fill(ch,uniq('channel')); fill(ct,uniq('customer_type'));
[c,ch,ct].forEach(x=>x.addEventListener('change',render)); render();
</script>
</body></html>
"""
    output.write_text(html, encoding="utf-8")
    return output


def build_notebook_en() -> Path:
    output = DELIVERABLES_DIR / "marketing_campaign_analysis_en.ipynb"
    nb = nbf.v4.new_notebook()

    md = """
# Marketing Campaign Analysis Notebook (EN)

This notebook covers:
1. Data loading and quality checks.
2. KPI computation.
3. Segmented analysis by campaign and channel.
4. Visual analysis and executive conclusions.
""".strip()

    code1 = """
from pathlib import Path
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt

base_dir = Path.cwd()
candidates = [
    base_dir / 'Marketing_Campaign_Data.csv',
    base_dir / 'Marketing Analysis' / 'Marketing_Campaign_Data.csv',
    base_dir.parent / 'Marketing_Campaign_Data.csv',
]

data_path = next((p for p in candidates if p.exists()), None)
if data_path is None:
    raise FileNotFoundError('Could not find Marketing_Campaign_Data.csv')

df = pd.read_csv(data_path)
df = df.rename(columns={
    'Interaction ID': 'interaction_id',
    'Campaign Type': 'campaign',
    'Channel': 'channel',
    'Customer Type': 'customer_type',
    'Converted (1=yes, 0=no)': 'converted',
    'Time on Site (seconds)': 'time_on_site',
    'Sales ($)': 'sales'
})

sns.set_theme(style='whitegrid')
df.head()
""".strip()

    code2 = """
kpis = {
    'interactions': len(df),
    'total_revenue': df['sales'].sum(),
    'conversion_rate': df['converted'].mean(),
    'avg_time': df['time_on_site'].mean(),
    'avg_ticket': df.loc[df['converted'] == 1, 'sales'].mean(),
}
kpis
""".strip()

    code3 = """
by_channel = df.groupby('channel', as_index=False).agg(
    interactions=('interaction_id', 'count'),
    conversion_rate=('converted', 'mean'),
    revenue=('sales', 'sum')
).sort_values('revenue', ascending=False)

plt.figure(figsize=(10,6))
sns.barplot(data=by_channel, x='channel', y='revenue')
plt.title('Revenue by Channel')
plt.show()
""".strip()

    code4 = """
by_combo = df.groupby(['campaign','channel'], as_index=False).agg(
    interactions=('interaction_id','count'),
    conversion_rate=('converted','mean'),
    revenue=('sales','sum')
).sort_values('revenue', ascending=False)

top = by_combo.iloc[0]
print('Top combination:', top['campaign'], '+', top['channel'])
print('Revenue:', round(top['revenue'], 2))
""".strip()

    nb.cells = [
        nbf.v4.new_markdown_cell(md),
        nbf.v4.new_code_cell(code1),
        nbf.v4.new_code_cell(code2),
        nbf.v4.new_code_cell(code3),
        nbf.v4.new_code_cell(code4),
    ]

    nbf.write(nb, output)
    return output


def build_presentation_en(kpis: dict, tables: dict) -> Path:
    output = DELIVERABLES_DIR / "Marketing_Campaign_Executive_Presentation_EN.pptx"
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Marketing Campaign Performance"
    slide.placeholders[1].text = "Executive briefing: campaign, channel, and customer insights"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "KPI Snapshot"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = f"Total interactions: {kpis['interactions']:,}"
    for line in [
        f"Total revenue: ${kpis['total_revenue']:,.2f}",
        f"Conversion rate: {kpis['conversion_rate'] * 100:.2f}%",
        f"Average time on site: {kpis['avg_time']:.1f} sec",
        f"Average ticket (converted): ${kpis['avg_ticket']:.2f}",
    ]:
        p = tf.add_paragraph(); p.text = line

    chart_map = {
        "Revenue by Channel and Campaign": CHARTS_DIR / "revenue_by_channel_campaign.png",
        "Conversion Heatmap": CHARTS_DIR / "conversion_heatmap.png",
        "Revenue by Customer Type": CHARTS_DIR / "customer_revenue.png",
    }
    for title, path in chart_map.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        if path.exists():
            slide.shapes.add_picture(str(path), PptInches(0.6), PptInches(1.2), width=PptInches(12.1))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommendations"
    tf = slide.shapes.placeholders[1].text_frame
    top = tables["by_combo"].iloc[0]
    tf.text = f"Scale investment in {top['campaign']} + {top['channel']} as the top revenue combination."
    for line in [
        "Run A/B tests by channel to optimize conversion.",
        "Use separate messaging for New vs Existing customers.",
        "Apply weekly KPI governance to rebalance budget.",
    ]:
        p = tf.add_paragraph(); p.text = line

    prs.save(output)
    return output


def build_board_ready_en(kpis: dict, tables: dict) -> Path:
    output = DELIVERABLES_DIR / "Marketing_Campaign_Executive_Presentation_Board_Ready_EN.pptx"
    prs = Presentation()

    top = tables["by_combo"].iloc[0]

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Executive Growth Plan"
    slide.placeholders[1].text = "Board-ready summary for campaign optimization"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Decision Requested Today"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Approve a 4-week pilot to increase investment in B + Email."
    for line in [
        "Objective: grow revenue while keeping conversion stable.",
        "Scope: creative, landing page and customer segmentation tests.",
        "Governance: weekly KPI review and decision rules.",
    ]:
        p = tf.add_paragraph(); p.text = line

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Expected Financial Impact"
    scen_chart = CHARTS_DIR / "escenarios_impacto_ingresos.png"
    if scen_chart.exists():
        slide.shapes.add_picture(str(scen_chart), Inches(0.8), Inches(1.2), width=Inches(7.2))

    txt = slide.shapes.add_textbox(Inches(8.1), Inches(1.3), Inches(4.8), Inches(4.6)).text_frame
    txt.text = "Executive readout"
    txt.paragraphs[0].font.bold = True
    txt.paragraphs[0].font.size = Pt(16)
    for line in [
        f"Top current combo: {top['campaign']} + {top['channel']}",
        f"Current total revenue: ${kpis['total_revenue']:,.0f}",
        "Recommend approving pilot under controlled KPI guardrails.",
    ]:
        p = txt.add_paragraph(); p.text = line; p.font.size = Pt(13)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "30-60-90 Execution Plan"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "0-30: launch pilot and set baseline"
    for line in [
        "31-60: run A/B tests and optimize winning variants",
        "61-90: scale proven tactics and formalize budget rebalance",
        "Owners: Marketing Performance, BI, and Sales",
    ]:
        p = tf.add_paragraph(); p.text = line

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risks and Mitigation"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Risk: audience fatigue in top channel"
    for line in [
        "Mitigation: rotate creatives weekly and cap frequency",
        "Risk: volume growth with lower lead quality",
        "Mitigation: enforce minimum conversion and ticket thresholds",
        "Risk: seasonality bias",
        "Mitigation: compare against a 4-week historical control",
    ]:
        p = tf.add_paragraph(); p.text = line

    prs.save(output)
    return output


def build_report_en(kpis: dict, tables: dict) -> Path:
    output = DELIVERABLES_DIR / "Marketing_Campaign_Executive_Report_EN.docx"
    doc = Document()
    doc.add_heading("Marketing Campaign Executive Report", level=0)
    doc.add_paragraph("English version generated from the campaign analysis dataset.")

    doc.add_heading("KPI Overview", level=1)
    t = doc.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "Metric"
    t.rows[0].cells[1].text = "Value"
    rows = [
        ("Total interactions", f"{kpis['interactions']:,}"),
        ("Total revenue", f"${kpis['total_revenue']:,.2f}"),
        ("Conversion rate", f"{kpis['conversion_rate'] * 100:.2f}%"),
        ("Avg time on site", f"{kpis['avg_time']:.2f} sec"),
        ("Avg ticket (converted)", f"${kpis['avg_ticket']:.2f}"),
    ]
    for metric, value in rows:
        c = t.add_row().cells
        c[0].text = metric
        c[1].text = value

    doc.add_heading("Key Findings", level=1)
    top = tables["by_combo"].iloc[0]
    findings = [
        f"Top revenue combination is {top['campaign']} + {top['channel']}.",
        "Email is a major revenue driver.",
        "Campaign B scales efficiently with strong volume.",
    ]
    for f in findings:
        doc.add_paragraph(f, style="List Bullet")

    doc.add_heading("Recommendations", level=1)
    for rec in [
        "Scale budget on the leading campaign-channel combination.",
        "Run channel-level A/B tests for creative and landing pages.",
        "Separate acquisition strategy for new and existing customers.",
        "Track KPI thresholds weekly and rebalance budget quickly.",
    ]:
        doc.add_paragraph(rec, style="List Number")

    for chart in [
        CHARTS_DIR / "revenue_by_channel_campaign.png",
        CHARTS_DIR / "conversion_heatmap.png",
        CHARTS_DIR / "customer_revenue.png",
    ]:
        if chart.exists():
            doc.add_picture(str(chart), width=Inches(6.4))

    doc.save(output)
    return output


def main() -> None:
    DELIVERABLES_DIR.mkdir(parents=True, exist_ok=True)
    df = load_data()
    kpis, tables = compute_tables(df)

    generated = [
        build_summary_en(kpis, tables),
        build_executive_dashboard_en(df, kpis, tables),
        build_interactive_dashboard_en(df),
        build_notebook_en(),
        build_presentation_en(kpis, tables),
        build_board_ready_en(kpis, tables),
        build_report_en(kpis, tables),
    ]

    print("Generated English versions:")
    for path in generated:
        print(f"- {path}")


if __name__ == "__main__":
    main()
